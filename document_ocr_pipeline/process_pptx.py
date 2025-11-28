#!/usr/bin/env python3
"""
PPTX 完整处理管道
生成与 PDF 流程一致的输出结构
"""

import sys
import json
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
from PIL import Image
import subprocess
import re

# 添加项目路径
sys.path.insert(0, str(Path(__file__).parent.parent))

from document_ocr_pipeline.extract_document import DocumentExtractor
from document_ocr_pipeline.visualize_extraction import visualize_extraction
from src.utils import get_soffice_command

# 尝试导入 VLM (可选依赖)
try:
    from src.models import VisionModel
    HAS_VLM = True
except ImportError:
    HAS_VLM = False
    print("⚠️  VLM 模块未安装，将跳过智能文本修正功能")


def detect_problem_content(text_blocks):
    """
    检测是否需要 VLM 修正
    
    返回: (需要修正, 原因, 统计信息)
    """
    if not text_blocks:
        return False, "无文本块", {}
    
    # 统计信息
    confidences = [block.get('confidence', 0) for block in text_blocks]
    avg_confidence = sum(confidences) / len(confidences) if confidences else 1.0
    low_conf_count = len([c for c in confidences if c < 0.7])
    low_conf_ratio = low_conf_count / len(confidences) if confidences else 0
    
    all_text = ' '.join([block.get('text', '') for block in text_blocks])
    
    # 检测特殊字符/乱码
    garbled_chars = len(re.findall(r'[□■�？?旬]', all_text))
    garbled_ratio = garbled_chars / max(len(all_text), 1)
    
    # 检测文件列表特征
    has_file_extensions = bool(re.search(r'\.(dmg|pkg|tar|gz|zip|app|png|jpg)', all_text, re.IGNORECASE))
    has_tree_symbols = any(char in all_text for char in ['三', '├', '└', '│', '─'])
    has_slash = '/' in all_text
    
    is_file_list = has_file_extensions and (has_tree_symbols or has_slash)
    
    # 检测短行多行特征（文件列表典型特征）
    lines = all_text.split('\n')
    short_lines = [line for line in lines if 0 < len(line.strip()) < 50]
    is_multi_short_lines = len(short_lines) > 5
    
    # 检测思维导图/关系图（树形符号密度）
    tree_symbols = sum(all_text.count(s) for s in ['├', '└', '│', '──', '─'])
    arrow_symbols = sum(all_text.count(s) for s in ['→', '←', '↓', '↑', '⇒', '⇐', '▶', '◀'])
    is_mindmap = (tree_symbols > 5 or arrow_symbols > 3) and len(text_blocks) > 8
    
    stats = {
        'avg_confidence': avg_confidence,
        'low_conf_ratio': low_conf_ratio,
        'garbled_ratio': garbled_ratio,
        'is_file_list': is_file_list,
        'is_multi_short_lines': is_multi_short_lines,
        'is_mindmap': is_mindmap,
        'tree_symbols_count': tree_symbols,
        'arrow_symbols_count': arrow_symbols
    }
    
    # 触发条件（满足任一）
    # 注释严格条件，采用更宽松的策略让 VLM 有机会介入修正
    # if avg_confidence < 0.5:  # 平均置信度阈值
    #     return True, f"平均置信度过低 ({avg_confidence:.2f})", stats
    # elif garbled_ratio > 0.03:  # 乱码字符阈值
    #     return True, f"检测到乱码字符 ({garbled_ratio:.1%})", stats
    
    # 宽松介入策略（60-80% 覆盖率）
    if avg_confidence < 0.8:  # 80% 以下就修正
        return True, f"识别质量可提升 (置信度 {avg_confidence:.2f})", stats
    elif garbled_ratio > 0.005:  # 0.5% 乱码即触发
        return True, f"检测到乱码 ({garbled_ratio:.1%})", stats
    elif stats.get('is_mindmap', False):  # 思维导图
        return True, "检测到思维导图/关系图", stats
    elif is_file_list or is_multi_short_lines:  # 特殊格式
        return True, "检测到列表结构", stats
    
    return False, "质量良好", stats


def refine_text_with_vlm(image_path, ocr_text, vlm_model, context_hint="", confidence_info=None):
    """
    使用 VLM 修正 OCR 文本
    
    Args:
        image_path: 图片路径
        ocr_text: OCR 原始文本
        vlm_model: VisionModel 实例
        context_hint: 上下文提示（如"文件列表"）
        confidence_info: 置信度信息 dict (avg_confidence, garbled_ratio)
    
    Returns:
        修正后的文本
    """
    if not HAS_VLM or not vlm_model:
        return ocr_text
    
    try:
        # 构建质量提示信息
        quality_note = ""
        correction_level = ""
        content_type_hint = ""
        if confidence_info:
            avg_conf = confidence_info.get('avg_confidence', 0)
            garbled = confidence_info.get('garbled_ratio', 0)
            is_mindmap = confidence_info.get('is_mindmap', False)
            is_file_list = confidence_info.get('is_file_list', False)
            
            if avg_conf < 0.5:
                quality_note = f"\n注意：OCR 识别质量较低（平均置信度 {avg_conf:.1%}），可能存在较多错误。"
                correction_level = "【激进修正模式】识别质量很低，需要大幅修正错别字和结构格式"
            elif avg_conf < 0.7:
                correction_level = "【中等修正模式】适度修正明显的错别字，保留大部分原文和结构格式"
            else:
                correction_level = "【保守修正模式】仅修正明显错误，保留格式和边距，保留原有结构格式"
            
            if garbled > 0.03:
                quality_note += f"\n注意：检测到 {garbled:.1%} 的乱码字符，请参考图片修正。"
            
            # 内容类型提示
            if is_mindmap:
                content_type_hint = "\n⚠️ **这是思维导图/关系图**，必须保留所有层级关系、分支结构、箭头方向！"
            elif is_file_list:
                content_type_hint = "\n⚠️ **这是文件列表/目录**，必须保留层级缩进和符号！"
        
        prompt = f"""请根据图片和 OCR 识别结果，修正以下文本中的错误：

OCR 原始结果：
{ocr_text}

识别质量信息：
{quality_note}
{correction_level}
{content_type_hint}

修正要求：
1. **错别字修正**（必须参考图片）：
   - 容器监控/应用监控/数据库监控 等IT术语
   - 常见错误：客器→容器、申间→空间、V志→日志、禺→域
   - 专有名词：CyberArk、Kong、API Gateway、CMDB

2. **格式保留**（禁止修改）：
   - 树形符号：├ │ └ ── 
   - 箭头符号：→ ← ↓ ↑ ⇒ ▶
   - 缩进层级：必须与原文一致
   - 换行位置：保持原有布局

3. **结构修复**（思维导图/关系图重点）：
   - **补充丢失的层级符号**（/, -, |, ├, └）
   - **恢复父子关系**（如 A → B → C 的流向）
   - **保持分支结构**（多个子节点必须全部展示）
   - 合并被错误分割的词语

4. **禁止行为**：
   - 不要添加原图中没有的内容
   - 不要改变节点之间的连接关系
   - 不要合并应该分开的分支
   - 不要删除看似重复但实际存在的内容

{f'提示：这是一个{context_hint}' if context_hint else ''}

请直接返回修正后的文本内容，不要有其他解释。"""

        result = vlm_model.extract_text_from_image(
            image_path=str(image_path),
            prompt=prompt
        )
        
        refined_text = result.get('text', ocr_text).strip()
        
        # 简单验证：如果修正后文本太短或太长（与原文相差10倍），可能有问题
        if len(refined_text) < len(ocr_text) * 0.3 or len(refined_text) > len(ocr_text) * 5:
            print(f"      ⚠️  VLM 修正结果异常，保持原文本")
            return ocr_text
        
        return refined_text
        
    except Exception as e:
        print(f"      ⚠️  VLM 修正失败: {e}")
        return ocr_text


def extract_slide_content(slide, slide_num, output_dir, ocr_engine='paddle'):
    """
    提取单页内容：文本 + 图片OCR + VLM组合
    """
    slide_data = {
        "page_number": slide_num,
        "statistics": {},
        "stage1_global": {},
        "stage3_vlm": {}
    }
    
    # ==================== 阶段 1: 文本直接提取 ====================
    print(f"\n  📝 阶段1: 提取文本内容（高优先级）...")
    
    extracted_text = {
        "title": "",
        "body": [],
        "notes": "",
        "tables": []
    }
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    placeholder = shape.placeholder_format
                    if placeholder.type == 1:  # Title
                        extracted_text["title"] = text
                    else:
                        extracted_text["body"].append(text)
                else:
                    extracted_text["body"].append(text)
        
        if shape.has_table:
            table = shape.table
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            extracted_text["tables"].append(table_data)
    
    if slide.has_notes_slide:
        extracted_text["notes"] = slide.notes_slide.notes_text_frame.text
    
    # 保存文本提取结果
    text_json_path = output_dir / f"page_{slide_num:03d}_extracted_text.json"
    with open(text_json_path, 'w', encoding='utf-8') as f:
        json.dump(extracted_text, f, ensure_ascii=False, indent=2)
    
    print(f"    ✓ 标题: {extracted_text['title'][:50] if extracted_text['title'] else '无'}")
    print(f"    ✓ 文本段落: {len(extracted_text['body'])} 个")
    print(f"    ✓ 表格: {len(extracted_text['tables'])} 个")
    
    # ==================== 阶段 1.5: 全页 OCR (用于获取坐标) ====================
    print(f"  👁️  阶段1.5: 全页 OCR (获取布局坐标)...")
    
    preview_image = f"page_{slide_num:03d}_300dpi.png"
    preview_path = output_dir / preview_image
    global_ocr_path = output_dir / f"page_{slide_num:03d}_global_ocr.json"
    visualized_image = f"page_{slide_num:03d}_visualized.png"
    visualized_path = output_dir / visualized_image
    
    if preview_path.exists():
        try:
            # 对全页预览图运行 OCR (300 DPI)
            extractor = DocumentExtractor(ocr_engine=ocr_engine)
            global_ocr_result = extractor.extract_from_image(str(preview_path))
            
            with open(global_ocr_path, 'w', encoding='utf-8') as f:
                json.dump(global_ocr_result, f, ensure_ascii=False, indent=2)
            
            # 生成可视化图
            visualize_extraction(str(preview_path), str(global_ocr_path), str(visualized_path))
            
            print(f"    ✓ 全页 OCR 完成: {len(global_ocr_result.get('text_blocks', []))} 个文本块")
            print(f"    ✓ 坐标数据已保存: {global_ocr_path.name}")
            
            # ==================== 阶段 1.6: 大字检测与 150 DPI 补充识别 ====================
            text_blocks = global_ocr_result.get('text_blocks', [])
            if text_blocks:
                # 获取图片尺寸
                with Image.open(preview_path) as img:
                    img_width, img_height = img.size
                    img_area = img_width * img_height
                
                # 计算整体置信度统计
                confidences = [block.get('confidence', 0) for block in text_blocks]
                avg_confidence = sum(confidences) / len(confidences) if confidences else 1.0
                low_conf_ratio = len([c for c in confidences if c < 0.7]) / len(confidences) if confidences else 0
                
                # 检测是否存在"疑似大字区域"
                large_blocks = []
                for block in text_blocks:
                    bbox = block.get('bbox', [0, 0, 0, 0])
                    if len(bbox) == 4:
                        x1, y1, x2, y2 = bbox
                        block_width = x2 - x1
                        block_height = y2 - y1
                        block_area = block_width * block_height
                        
                        # 判断条件：单个文字块面积 > 图片面积的 10%，或尺寸 > 300x300 px
                        if block_area > img_area * 0.1 or (block_width > 300 and block_height > 300):
                            large_blocks.append(block)
                
                # 触发 150 DPI 缩小识别的条件（二选一）：
                # 1. 检测到大字块
                # 2. 整体置信度低（平均 < 0.65 或 超过 50% 的块 < 0.7）
                should_try_150dpi = (
                    len(large_blocks) > 0 or 
                    avg_confidence < 0.65 or 
                    low_conf_ratio > 0.5
                )
                
                if should_try_150dpi:
                    if large_blocks:
                        print(f"    🔍 检测到 {len(large_blocks)} 个大字区域，尝试 150 DPI 缩小识别...")
                    else:
                        print(f"    🔍 整体置信度较低 (平均: {avg_confidence:.2f}, 低置信度占比: {low_conf_ratio:.1%})，尝试 150 DPI 缩小识别...")
                    
                    # 生成 150 DPI 缩小版图片（缩小到原来的 50%）
                    preview_150dpi_path = output_dir / f"page_{slide_num:03d}_preview_150dpi.png"
                    with Image.open(preview_path) as img:
                        new_size = (img_width // 2, img_height // 2)
                        img_150dpi = img.resize(new_size, Image.LANCZOS)
                        img_150dpi.save(preview_150dpi_path)
                    
                    # 对 150 DPI 图片运行 OCR
                    ocr_150dpi_result = extractor.extract_from_image(str(preview_150dpi_path))
                    
                    # 将 150 DPI 的坐标还原到 300 DPI（坐标 x2）
                    for block in ocr_150dpi_result.get('text_blocks', []):
                        if 'bbox' in block and len(block['bbox']) == 4:
                            block['bbox'] = [coord * 2 for coord in block['bbox']]
                    
                    # 合并结果：优先使用高置信度的
                    merged_blocks = []
                    used_150dpi_indices = set()
                    
                    for block_300 in text_blocks:
                        bbox_300 = block_300.get('bbox', [0, 0, 0, 0])
                        best_match = block_300
                        
                        # 检查是否有 150 DPI 的结果覆盖同一区域且置信度更高
                        for idx, block_150 in enumerate(ocr_150dpi_result.get('text_blocks', [])):
                            if idx in used_150dpi_indices:
                                continue
                            
                            bbox_150 = block_150.get('bbox', [0, 0, 0, 0])
                            
                            # 判断两个框是否重叠（IoU > 0.3）
                            x1_300, y1_300, x2_300, y2_300 = bbox_300
                            x1_150, y1_150, x2_150, y2_150 = bbox_150
                            
                            x_overlap = max(0, min(x2_300, x2_150) - max(x1_300, x1_150))
                            y_overlap = max(0, min(y2_300, y2_150) - max(y1_300, y1_150))
                            overlap_area = x_overlap * y_overlap
                            
                            area_300 = (x2_300 - x1_300) * (y2_300 - y1_300)
                            area_150 = (x2_150 - x1_150) * (y2_150 - y1_150)
                            union_area = area_300 + area_150 - overlap_area
                            
                            if union_area > 0:
                                iou = overlap_area / union_area
                                if iou > 0.3:  # 重叠度 > 30%
                                    # 优先使用高置信度的结果
                                    conf_300 = block_300.get('confidence', 0)
                                    conf_150 = block_150.get('confidence', 0)
                                    
                                    if conf_150 > conf_300:
                                        best_match = block_150
                                        used_150dpi_indices.add(idx)
                                    break
                        
                        merged_blocks.append(best_match)
                    
                    # 添加未匹配的 150 DPI 结果
                    for idx, block_150 in enumerate(ocr_150dpi_result.get('text_blocks', [])):
                        if idx not in used_150dpi_indices:
                            merged_blocks.append(block_150)
                    
                    # 更新结果
                    improvement_count = len([b for b in merged_blocks if b.get('confidence', 0) > 0.9])
                    original_high_conf = len([b for b in text_blocks if b.get('confidence', 0) > 0.9])
                    
                    if improvement_count > original_high_conf:
                        global_ocr_result['text_blocks'] = merged_blocks
                        with open(global_ocr_path, 'w', encoding='utf-8') as f:
                            json.dump(global_ocr_result, f, ensure_ascii=False, indent=2)
                        
                        # 重新生成可视化
                        visualize_extraction(str(preview_path), str(global_ocr_path), str(visualized_path))
                        
                        print(f"    ✓ 150 DPI 补充识别完成: 合并后 {len(merged_blocks)} 个文本块 (高置信度: {original_high_conf} → {improvement_count})")
                    else:
                        print(f"    ℹ️  150 DPI 识别未带来明显改善，保持原结果")
                
        except Exception as e:
            print(f"    ⚠️ 全页 OCR 失败: {e}")
            # 如果失败，创建一个空的 OCR 结果以防报错
            with open(global_ocr_path, 'w', encoding='utf-8') as f:
                json.dump({"text_blocks": []}, f)
    else:
        print(f"    ⚠️ 预览图不存在，跳过全页 OCR: {preview_path.name}")
    
    # ==================== 阶段 2: 处理嵌入图片（OCR） ====================
    print(f"  🖼️  阶段2: 处理嵌入图片内容...")
    
    images = []
    image_ocr_results = []
    
    for idx, shape in enumerate(slide.shapes, 1):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_filename = f"page_{slide_num:03d}_img_{idx}.{image.ext}"
            image_path = output_dir / image_filename
            
            # 保存图片
            with open(image_path, "wb") as f:
                f.write(image.blob)
            
            try:
                with Image.open(io.BytesIO(image.blob)) as img:
                    width, height = img.size
                
                images.append({
                    "id": idx,
                    "path": str(image_path),
                    "format": image.ext,
                    "size": [width, height]
                })
                
                print(f"    ✓ 嵌入图片 {idx}: {width}x{height} ({image.ext})")
                
                # 对图片运行 OCR
                img_ocr_json_path = output_dir / f"page_{slide_num:03d}_img_{idx}_ocr.json"
                extractor = DocumentExtractor(ocr_engine=ocr_engine)
                ocr_result = extractor.extract_from_image(str(image_path))
                
                with open(img_ocr_json_path, 'w', encoding='utf-8') as f:
                    json.dump(ocr_result, f, ensure_ascii=False, indent=2)
                
                # 生成可视化
                img_vis_path = output_dir / f"page_{slide_num:03d}_img_{idx}_visualized.png"
                visualize_extraction(str(image_path), str(img_ocr_json_path), str(img_vis_path))
                
                image_ocr_results.append({
                    "image_id": idx,
                    "ocr_json": str(img_ocr_json_path),
                    "visualized": str(img_vis_path),
                    "text_blocks_count": len(ocr_result.get('text_blocks', []))
                })
                
                print(f"      ✓ OCR: {len(ocr_result.get('text_blocks', []))} 个文本块")
            except Exception as e:
                print(f"      ✗ 图片处理失败: {e}")
    
    # ==================== 阶段 3: VLM 处理（带文本上下文） ====================
    print(f"  🤖 阶段3: VLM 综合分析...")
    
    # 构建 VLM 提示（包含已提取的文本）
    vlm_context = {
        "extracted_text": extracted_text,
        "images": images,
        "image_ocr_results": image_ocr_results
    }
    
    # 生成 VLM 输入提示
    vlm_prompt = f"""# 幻灯片第 {slide_num} 页综合分析

## 已提取的文本内容（高可信度，来自PPT原始数据）：

### 标题
{extracted_text['title'] or '无'}

### 正文内容
"""
    for i, body in enumerate(extracted_text['body'], 1):
        vlm_prompt += f"{i}. {body}\n"
    
    if extracted_text['tables']:
        vlm_prompt += "\n### 表格\n"
        for t_idx, table in enumerate(extracted_text['tables'], 1):
            vlm_prompt += f"表格 {t_idx}:\n"
            for row in table[:3]:
                vlm_prompt += f"  {' | '.join(row)}\n"
    
    if extracted_text['notes']:
        vlm_prompt += f"\n### 备注\n{extracted_text['notes']}\n"
    
    vlm_prompt += f"""

## 图片OCR结果（{len(image_ocr_results)} 张图片）：
"""
    for ocr_res in image_ocr_results:
        vlm_prompt += f"- 图片 {ocr_res['image_id']}: {ocr_res['text_blocks_count']} 个文本块\n"
    
    vlm_prompt += """

## VLM 任务：
请综合上述信息，生成完整的页面内容描述：
1. 确认并整合已提取的文本
2. 补充图片中的额外信息（图表、示意图、图标等）
3. 识别页面的整体类型和主题
4. 提取关键信息和结构

返回 JSON 格式。
"""
    
    vlm_prompt_path = output_dir / f"page_{slide_num:03d}_vlm_prompt.txt"
    with open(vlm_prompt_path, 'w', encoding='utf-8') as f:
        f.write(vlm_prompt)
    
    # 这里可以调用 VLM（如果有的话）
    # 暂时保存提示和上下文
    vlm_context_path = output_dir / f"page_{slide_num:03d}_vlm_context.json"
    with open(vlm_context_path, 'w', encoding='utf-8') as f:
        json.dump(vlm_context, f, ensure_ascii=False, indent=2)
    
    print(f"    ✓ VLM上下文和提示已保存")
    
    # ==================== 构建最终页面数据 ====================
    # 整合所有文本
    all_text = []
    if extracted_text['title']:
        all_text.append(extracted_text['title'])
    all_text.extend(extracted_text['body'])
    if extracted_text['notes']:
        all_text.append(f"备注: {extracted_text['notes']}")
    
    # 添加表格文本
    for table in extracted_text['tables']:
        for row in table:
            all_text.append(' | '.join(row))
    
    # ==================== 阶段 3.5: VLM 智能文本修正（按需触发） ====================
    # 添加图片OCR文本（只保留高置信度）
    MIN_CONFIDENCE = 0.15
    vlm_model = None
    
    for ocr_res in image_ocr_results:
        try:
            with open(ocr_res['ocr_json'], 'r', encoding='utf-8') as f:
                ocr_data = json.load(f)
                # 从 text_blocks 提取高置信度文本
                if ocr_data.get('text_blocks'):
                    high_confidence_blocks = [
                        block for block in ocr_data['text_blocks']
                        if block.get('confidence', 0) >= MIN_CONFIDENCE and block.get('text')
                    ]
                    
                    # 检测是否需要 VLM 修正
                    all_blocks = ocr_data['text_blocks']
                    needs_refinement, reason, stats = detect_problem_content(all_blocks)
                    
                    if needs_refinement and high_confidence_blocks:
                        print(f"      🔍 触发 VLM 修正 - {reason}")
                        print(f"         (平均置信度: {stats['avg_confidence']:.2f}, 乱码率: {stats['garbled_ratio']:.1%})")
                        
                        # 延迟初始化 VLM（只有需要时才加载）
                        if vlm_model is None and HAS_VLM:
                            try:
                                vlm_model = VisionModel()
                                print(f"      ✓ VLM 模型已加载")
                            except Exception as e:
                                print(f"      ⚠️  VLM 初始化失败: {e}")
                        
                        if vlm_model:
                            # 获取原始图片路径
                            img_id = ocr_res['image_id']
                            img_path = None
                            for img_info in images:
                                if img_info['id'] == img_id:
                                    img_path = img_info['path']
                                    break
                            
                            if img_path:
                                # 原始 OCR 文本
                                original_text = ' '.join([block['text'] for block in high_confidence_blocks])
                                
                                # VLM 修正
                                context_hint = "文件列表" if stats.get('is_file_list') else ""
                                confidence_info = {
                                    'avg_confidence': stats['avg_confidence'],
                                    'garbled_ratio': stats['garbled_ratio'],
                                    'is_mindmap': stats.get('is_mindmap', False),
                                    'is_file_list': stats.get('is_file_list', False)
                                }
                                refined_text = refine_text_with_vlm(
                                    image_path=img_path,
                                    ocr_text=original_text,
                                    vlm_model=vlm_model,
                                    context_hint=context_hint,
                                    confidence_info=confidence_info
                                )
                                
                                if refined_text != original_text:
                                    print(f"      ✓ VLM 修正完成 ({len(original_text)} → {len(refined_text)} 字符)")
                                    all_text.append(f"[图片 {ocr_res['image_id']}-VLM修正] {refined_text}")
                                else:
                                    all_text.append(f"[图片 {ocr_res['image_id']}-高置信度] " + original_text)
                            else:
                                # 找不到图片路径，使用原始文本
                                image_texts = [block['text'] for block in high_confidence_blocks]
                                all_text.append(f"[图片 {ocr_res['image_id']}-高置信度] " + ' '.join(image_texts))
                        else:
                            # VLM 不可用，使用原始文本
                            image_texts = [block['text'] for block in high_confidence_blocks]
                            all_text.append(f"[图片 {ocr_res['image_id']}-高置信度] " + ' '.join(image_texts))
                    
                    elif high_confidence_blocks:
                        # 质量良好，直接使用 OCR 文本
                        image_texts = [block['text'] for block in high_confidence_blocks]
                        all_text.append(f"[图片 {ocr_res['image_id']}-高置信度] " + ' '.join(image_texts))
                    
                    # 记录过滤情况
                    low_conf_count = len(ocr_data['text_blocks']) - len(high_confidence_blocks)
                    if low_conf_count > 0:
                        print(f"      ℹ️  图片 {ocr_res['image_id']}: 过滤了 {low_conf_count} 个低置信度文本块")
        except Exception as e:
            print(f"      ⚠️  无法读取图片OCR结果: {e}")
    
    combined_text = '\n\n'.join(all_text)
    
    # 计算整体OCR置信度（包括全局OCR和图片OCR）
    all_confidences = []
    
    # 从全局OCR获取置信度
    try:
        global_ocr_path = output_dir / f"page_{slide_num:03d}_global_ocr.json"
        if global_ocr_path.exists():
            with open(global_ocr_path, 'r', encoding='utf-8') as f:
                global_ocr = json.load(f)
                for block in global_ocr.get('text_blocks', []):
                    conf = block.get('confidence', 0)
                    if conf > 0:
                        all_confidences.append(conf)
    except Exception:
        pass
    
    # 从图片OCR获取置信度
    for ocr_res in image_ocr_results:
        try:
            with open(ocr_res['ocr_json'], 'r', encoding='utf-8') as f:
                ocr_data = json.load(f)
                for block in ocr_data.get('text_blocks', []):
                    conf = block.get('confidence', 0)
                    if conf > 0:
                        all_confidences.append(conf)
        except Exception:
            pass
    
    avg_ocr_confidence = sum(all_confidences) / len(all_confidences) if all_confidences else 0.0
    
    # 统计信息
    slide_data['statistics'] = {
        "total_text_blocks": len(all_text),
        "total_images": len(images),
        "has_title": bool(extracted_text['title']),
        "has_tables": len(extracted_text['tables']) > 0,
        "has_notes": bool(extracted_text['notes']),
        "avg_ocr_confidence": round(avg_ocr_confidence, 3)  # 添加平均置信度
    }
    
    # Stage1 信息（模拟 PDF 的结构）
    slide_data['stage1_global'] = {
        "image": preview_image,
        "ocr_json": f"page_{slide_num:03d}_global_ocr.json", # 指向包含坐标的 OCR 结果
        "text_source": "direct_extraction_plus_ocr"
    }
    
    # Stage2 OCR 可视化信息
    # 注意：visualized_path 已经在阶段 1.5 中通过 visualize_extraction 生成
    # 不要覆盖它，否则绿色 OCR 框会丢失
    
    if not visualized_path.exists():
        # 只有在生成失败时才降级处理
        if preview_path.exists():
            import shutil
            shutil.copy2(preview_path, visualized_path)
        elif images and image_ocr_results:
            # 降级：使用第一张图片的 OCR 可视化
            first_vis = next((r['visualized'] for r in image_ocr_results if r['image_id'] == images[0]['id']), None)
            if first_vis and Path(first_vis).exists():
                import shutil
                shutil.copy2(first_vis, visualized_path)
    
    # Stage3 VLM 信息
    slide_data['stage3_vlm'] = {
        "vlm_prompt": str(vlm_prompt_path.name),
        "vlm_context": str(vlm_context_path.name),
        "text_combined": combined_text
    }
    
    return slide_data, combined_text


def process_pptx(pptx_path, output_dir, ocr_engine='paddle'):
    """
    完整处理 PPTX 文件
    生成与 adaptive_ocr_pipeline.py 相同的输出结构
    """
    print(f"🚀 开始处理 PPTX: {pptx_path}")
    print(f"📂 输出目录: {output_dir}")
    print(f"🔧 OCR引擎: {ocr_engine}")
    
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = Path(pptx_path)
    
    # ==================== 步骤 0: 使用 LibreOffice 转换为 PDF 并渲染预览图 ====================
    print(f"\n{'='*70}")
    print(f"📄 步骤 0: 生成页面预览图（LibreOffice 渲染）")
    print(f"{'='*70}")
    
    temp_pdf = output_dir / f"{pptx_path.stem}_temp.pdf"
    
    # 获取 LibreOffice 命令
    soffice_cmd = get_soffice_command()
    if not soffice_cmd:
        print("  ⚠️  警告: 未找到 LibreOffice (soffice)，跳过预览图生成")
        print("  提示: 安装 LibreOffice 并确保 soffice 命令在 PATH 中")
        print("  macOS: brew install --cask libreoffice")
        total_slides = None
    else:
        try:
            # 调用 LibreOffice 转换 PPTX -> PDF
            print(f"  ⏳ 转换 PPTX 为 PDF (使用: {soffice_cmd})...")
            subprocess.run([
                soffice_cmd,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(output_dir),
                str(pptx_path)
            ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            
            # LibreOffice 输出的 PDF 文件名与输入文件名相同（仅扩展名不同）
            generated_pdf = output_dir / f"{pptx_path.stem}.pdf"
            if generated_pdf.exists() and generated_pdf != temp_pdf:
                generated_pdf.rename(temp_pdf)
            
            print(f"  ✓ PDF 已生成: {temp_pdf.name}")
            
            # 使用 pdfplumber 渲染每一页为图片
            import pdfplumber
            import cv2
            import numpy as np
            
            with pdfplumber.open(temp_pdf) as pdf:
                total_slides = len(pdf.pages)
                print(f"  📄 PDF 页数: {total_slides}")
                
                for page_num, page in enumerate(pdf.pages, 1):
                    # 渲染为高质量图片（300 DPI）
                    img = page.to_image(resolution=300)
                    img_array = np.array(img.original)
                    
                    # 保存为 page_XXX_300dpi.png（与 PDF 流程命名一致）
                    preview_path = output_dir / f"page_{page_num:03d}_300dpi.png"
                    cv2.imwrite(str(preview_path), cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR))
                    
                    height, width = img_array.shape[:2]
                    print(f"  ✓ 第 {page_num} 页: {width}x{height}px -> {preview_path.name}")
            
            # 删除临时 PDF 文件
            temp_pdf.unlink()
            print(f"  ✓ 预览图生成完成，临时 PDF 已清理")
            
        except Exception as e:
            print(f"  ⚠️  预览图生成失败: {e}")
            total_slides = None
    
    # ==================== 继续原有的内容提取流程 ====================
    prs = Presentation(str(pptx_path))
    if total_slides is None:
        total_slides = len(prs.slides)
    
    print(f"\n📄 总页数: {total_slides}")
    
    result = {
        "source_file": str(pptx_path),
        "file_type": "pptx",
        "total_pages": total_slides,
        "ocr_engine": ocr_engine,
        "pages": []
    }
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"\n{'='*70}")
        print(f"📄 处理第 {slide_idx}/{total_slides} 页")
        print(f"{'='*70}")
        
        slide_data, combined_text = extract_slide_content(
            slide, slide_idx, output_dir, ocr_engine
        )
        
        result["pages"].append(slide_data)
        
        print(f"\n✅ 第 {slide_idx} 页完成")
        print(f"  - 文本块: {slide_data['statistics']['total_text_blocks']}")
        print(f"  - 图片: {slide_data['statistics']['total_images']}")
        print(f"  - 字符数: {len(combined_text)}")
    
    # 保存完整结果（与 PDF 的 complete_adaptive_ocr.json 格式一致）
    complete_json = output_dir / "complete_adaptive_ocr.json"
    with open(complete_json, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    
    # 生成 complete_document.json（用于 ES 索引，标准格式）
    pages_for_index = []
    for page in result['pages']:
        page_num = page['page_number']
        stage1 = page.get('stage1_global', {})
        stage3 = page.get('stage3_vlm', {})
        stats = page.get('statistics', {})
        
        # 获取图片路径
        image_filename = stage1.get('image', f'page_{page_num:03d}_300dpi.png')
        image_path = output_dir / image_filename
        
        # 获取文本内容
        text_combined = stage3.get('text_combined', '')
        
        pages_for_index.append({
            'page_number': page_num,
            'image_path': str(image_path),
            'image_filename': image_filename,
            'content': {
                'full_text_cleaned': text_combined,
                'full_text_raw': text_combined,
                'key_fields': [],
                'tables': []
            },
            'ocr_data': {
                'text_blocks': []
            },
            'metadata': {
                'extraction_method': 'pptx_ocr_pipeline',
                'ocr_engine': ocr_engine,
                'avg_ocr_confidence': stats.get('avg_ocr_confidence', 0.0),
                'vlm_refined': False
            }
        })
    
    complete_document_path = output_dir / "complete_document.json"
    with open(complete_document_path, 'w', encoding='utf-8') as f:
        json.dump({'pages': pages_for_index}, f, ensure_ascii=False, indent=2)
    
    print(f"\n{'='*70}")
    print(f"✅ 处理完成！")
    print(f"{'='*70}")
    print(f"📊 统计:")
    print(f"  - 总页数: {total_slides}")
    print(f"  - 输出文件: {complete_json}")
    print(f"  - 索引文件: {complete_document_path}")
    print(f"  - 输出目录: {output_dir.absolute()}")
    
    return result


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPTX 完整处理管道")
    parser.add_argument("input_file", help="输入 PPTX 文件路径")
    parser.add_argument("-o", "--output", default=None, help="输出目录（默认：pptx_output）")
    parser.add_argument("--ocr-engine", choices=['easy', 'paddle', 'vision'], 
                       default='paddle', help="OCR引擎选择")
    
    args = parser.parse_args()
    
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"❌ 文件不存在: {input_path}")
        sys.exit(1)
    
    if args.output:
        output_dir = Path(args.output)
    else:
        # 默认输出目录：文件名_adaptive
        output_dir = Path(input_path.stem.replace(' ', '_') + "_adaptive")
    
    try:
        result = process_pptx(input_path, output_dir, args.ocr_engine)
        print(f"\n🎉 成功！可以使用此输出目录集成到系统中。")
    except Exception as e:
        print(f"\n❌ 处理失败: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
